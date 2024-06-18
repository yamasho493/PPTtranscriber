from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def get_new_and_existing_month_club_performance(prs, slide_number, sheet, row_number):
    slide = prs.slides[slide_number]
    performance_list = []
    for i in range(4):
        performance_list.append(sheet.iloc[row_number + i, 1:5])

    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            for list_index, list in enumerate(performance_list):
                for value_index, value in enumerate(list):
                    cell = table.cell(list_index + 1, value_index)
                    tf = cell.text_frame
                    if isinstance(value, int):
                        formatted_number = (f"{round(value, -3):,}").rstrip("0").rstrip(",")
                        cell.text = str(formatted_number)
                        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                    elif isinstance(value, float):
                        formatted_number = f"{value:.1%}"
                        cell.text = str(formatted_number)
                        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                    else:
                        cell.text = value
                        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
                    tf.paragraphs[0].font.name = "Meiryo UI"
                    tf.paragraphs[0].font.size = Pt(32)

def get_customer_number_and_unitprice_by_category(prs, slide_number, sheet, row_number):
    slide = prs.slides[slide_number]
    performance_list = []
    for i in range(8):
        performance_list.append(sheet.iloc[row_number + i, 1:11])
    
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            for list_index, list in enumerate(performance_list):
                for value_index, value in enumerate(list):
                    cell = table.cell(list_index + 2, value_index)
                    tf = cell.text_frame
                    if isinstance(value, int):
                        formatted_number = f"{value:,}"
                        cell.text = str(formatted_number)
                        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                    elif isinstance(value, float) and value < 1:
                        formatted_number = f"{value:.1%}"
                        cell.text = str(formatted_number)
                        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                        tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                    elif isinstance(value, float) and value >= 1:
                        formatted_number = f"{value:.1%}"
                        cell.text = str(formatted_number)
                        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                    else:
                        cell.text = value
                        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
                    tf.paragraphs[0].font.name = "Meiryo UI"
                    tf.paragraphs[0].font.size = Pt(14)

def get_monthly_regional_performance(prs, slide_number, sheet):
    slide = prs.slides[slide_number]
    perfomance_list = []

    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            first_row_cells = [cell.text for cell in table.rows[0].cells]
            slice_first_row_cells = first_row_cells[1:]
            for cell in slice_first_row_cells:
                regional_performance_list = []
                indices = (sheet == cell).values.nonzero()
                first_match_row = indices[0][0]
                first_match_colunm = indices[1][0]
                for i in range(3):
                    regional_performance_list.append(sheet.iloc[first_match_row + i + 1, first_match_colunm])
                perfomance_list.append(regional_performance_list)
                
            for list_index, list in enumerate(perfomance_list):
                for value_index, value in enumerate(list):
                    cell = table.cell(value_index + 1, list_index + 1)
                    tf = cell.text_frame
                    formatted_number = f"{value:.1%}"
                    cell.text = str(formatted_number)
                    tf.paragraphs[0].font.name = "Meiryo UI"
                    tf.paragraphs[0].font.size = Pt(16)
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def insert_word_to_column(prs, slide_number, word):
    slide = prs.slides[slide_number]

    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            cell = table.cell(0, 1)
            tf = cell.text_frame
            cell.text = word
            tf.paragraphs[0].font.name = "Meiryo UI"
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER


if __name__ =="__main__":
    get_new_and_existing_month_club_performance()
    get_customer_number_and_unitprice_by_category()
    get_monthly_regional_performance()
    insert_word_to_column()
