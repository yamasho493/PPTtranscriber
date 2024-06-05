from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def get_month_performances(prs, slide_number, sheet):
    slide = prs.slides[slide_number]
    month_performance_value_list = sheet.iloc[6:9, 2]

    for index, value in enumerate(month_performance_value_list):
        left = Inches(3.7)
        top = Inches(2.1 + 0.59 * index)
        height = Inches(0.68)
        width = Inches(2.76)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        paragraph = tf.paragraphs[0]
        paragraph.text = value
        paragraph.font.name = "Meiryo UI"
        paragraph.font.size = Pt(32)

def get_round_rates(prs, slide_number, sheet):
    slide = prs.slides[slide_number]
    early_period_rate_list = sheet.iloc[6:9, 4]

    for index, value in enumerate(early_period_rate_list):
        left = Inches(10)
        top = Inches(2.1 + 0.59 * index)
        height = Inches(0.68)
        width = Inches(2.76)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        paragraph = tf.paragraphs[0]
        formatted_value = f"{value:.1%}"
        paragraph.text = str(formatted_value)
        paragraph.font.name = "Meiryo UI"
        paragraph.font.size = Pt(32)

def get_club_customer_number_rate(prs, slide_number, sheet, row_number1, row_number2):
    slide = prs.slides[slide_number]
    club_list = []
    for i in range(3):
        club_list.append([sheet.iloc[row_number1, 2 + i * 2], sheet.iloc[row_number2, 2 + i * 2]])

    for list_index, list in enumerate(club_list):
        for value_index, value in enumerate(list):
            left = Inches(2.35 + 3.4 * list_index)
            top = Inches(3.65 + 2.3 * value_index)
            height = Inches(2.43)
            width = Inches(1.142)
            square = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, height, width)
            tf = square.text_frame
            paragraph = tf.paragraphs[0]
            formatted_value = f"{value:.1%}"
            paragraph.text = str(formatted_value)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            paragraph.font.name = "Meiryo UI"
            paragraph.font.size = Pt(36)
            #影
            square.fill.background()
            line = square.line
            line.color.rgb = RGBColor(0, 0, 0)
            line.width = Pt(0.75)
            # shadow = square.shadow
            # shadow.inherit = False
            # shadow.type = MSO_SHADOW.OUTER
            # shadow.blur_radius = Inches(1)
            # shadow.distance = Inches(1)

def get_new_development(prs, slide_number, sheet, row_number1, row_number2, column_number, misalignment):
    slide = prs.slides[slide_number]
    performance_list = sheet.iloc[row_number1:row_number2, column_number]
    for index, value in enumerate(performance_list):
        left = Inches(4.15)
        top = Inches(2.8 + 0.5 * index + 2.638 * misalignment)
        height = Inches(0.68)
        width = Inches(3.23)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        paragraph = tf.paragraphs[0]
        paragraph.text = value
        paragraph.font.name = "Meiryo UI"
        paragraph.font.size = Pt(32)
        paragraph.alignment = PP_ALIGN.RIGHT

def get_new_development_early_period_rate(prs, slide_number, sheet, row_number1, row_number2, column_number, misalignment):
    slide = prs.slides[slide_number]
    early_period_rate_list = sheet.iloc[row_number1:row_number2, column_number]
    for index, value in enumerate(early_period_rate_list):
        left = Inches(10.18)
        top = Inches(2.8 + 0.5 * index + 2.638 * misalignment)
        height = Inches(0.68)
        width = Inches(2.19)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        paragraph = tf.paragraphs[0]
        if isinstance(value, int) or isinstance(value, float):
            formatted_value = f"{value:.1%}"
            paragraph.text = str(formatted_value)
            paragraph.alignment = PP_ALIGN.RIGHT
        else:
            paragraph.text = str(value)
            paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = "Meiryo UI"
        paragraph.font.size = Pt(32)

if __name__ == "__main__":
    get_month_performances()
    get_round_rates()
    get_club_customer_number_rate()
    get_new_development()
    get_new_development_early_period_rate()
