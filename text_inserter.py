from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def insert_ppt_titles(prs, slide_number, ppt_titles):
    slide = prs.slides[slide_number]
    title_shape = slide.shapes.title
    title_shape.left = Inches(1.66)
    title_shape.top = Inches(1.49)
    title_tf = title_shape.text_frame
    title_tf.add_paragraph()
    for index, title in enumerate(ppt_titles):
        paragraph = title_tf.paragraphs[index]
        paragraph.text = title
        if index == 0:
            paragraph.font.size = Pt(44)
        elif index == 1:
            paragraph.font.size = Pt(25)
        else:
            return ""
        paragraph.font.name = "Meiryo UI"
        paragraph.font.color.rgb = RGBColor(0, 32, 96)
        paragraph.font.bold = True

def insert_slide_titie(prs, slide_number, title_text):
    slide = prs.slides[slide_number]
    title_shape = slide.shapes.title
    title_tf = title_shape.text_frame
    paragraph = title_tf.paragraphs[0]
    paragraph.text = title_text
    paragraph.font.name = "Meiryo UI"
    paragraph.font.color.rgb = RGBColor(0, 32, 96)
    paragraph.font.size = Pt(32)
    paragraph.font.bold = True

def insert_sub_title(prs, slide_number, sub_title, top, left):
    slide = prs.slides[slide_number]
    left = Inches(left)
    top = Inches(top)
    height = Inches(0.64)
    width = Inches(5.76)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    paragraph = tf.paragraphs[0]
    paragraph.text = sub_title
    paragraph.font.name = "Meiryo UI"
    paragraph.font.size = Pt(32)

if __name__ == "__main__":
    insert_ppt_titles()
    insert_slide_titie()
    insert_sub_title()
